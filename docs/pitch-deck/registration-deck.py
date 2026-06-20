#!/usr/bin/env python3
"""Registration pitch deck for OneVoice AI Challenge — Kami-inspired PPTX."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import argparse

# Design tokens
PARCHMENT   = RGBColor(0xf5, 0xf4, 0xed)
IVORY       = RGBColor(0xfa, 0xf9, 0xf5)
BRAND       = RGBColor(0x1B, 0x36, 0x5D)
NEAR_BLACK  = RGBColor(0x14, 0x14, 0x13)
DARK_WARM   = RGBColor(0x3d, 0x3d, 0x3a)
OLIVE       = RGBColor(0x50, 0x4e, 0x49)
STONE       = RGBColor(0x6b, 0x6a, 0x64)
BORDER      = RGBColor(0xe8, 0xe6, 0xdc)
WHITE       = RGBColor(0xff, 0xff, 0xff)

SERIF  = "XCharter"  # Charter extended; TeX Live OTF on Linux, 
                   # falls to Charter on macOS, Georgia on Windows
MONO   = "JetBrains Mono"
SANS   = SERIF

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Helpers ───────────────────────────────────────────────

def blank_slide(prs, bg=PARCHMENT):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 0, 0, prs.slide_width, prs.slide_height)
    bg_shp.fill.solid()
    bg_shp.fill.fore_color.rgb = bg
    bg_shp.line.fill.background()
    try:
        bg_shp.shadow.inherit = False
    except Exception:
        pass
    return s


def textbox(slide, text, left, top, width, height,
            font=SANS, size=18, bold=False, color=NEAR_BLACK,
            align=PP_ALIGN.LEFT, vanchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = vanchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def line(slide, left, top, width, color=BRAND, weight=1):
    shape = slide.shapes.add_connector(1, left, top, left + width, top)
    shape.line.color.rgb = color
    shape.line.width = Pt(int(weight))
    return shape


def card(slide, left, top, width, height, fill=IVORY, border=BORDER, border_pt=0.5):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = border
    s.line.width = Pt(border_pt)
    try:
        s.shadow.inherit = False
    except Exception:
        pass
    return s


def multi_text(slide, lines, left, top, width, height,
               font=SANS, size=18, color=DARK_WARM,
               line_spacing=1.5, bullet=False):
    """Multi-line text block — each line is a separate paragraph."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0

    for i, line_text in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            from pptx.oxml.ns import qn
            new_p = tf._p_dangerous  # internal — add paragraph
            # simpler: use add_paragraph from text_frame
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(4)
        p.space_before = Pt(0)
        prefix = "\u2022 " if bullet else ""
        run = p.add_run()
        run.text = prefix + line_text
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tb


# ── Slide templates ───────────────────────────────────────

def cover(prs, team_name, project, subtitle, date):
    s = blank_slide(prs)
    textbox(s, team_name,
            Inches(1), Inches(0.5), Inches(11.33), Inches(0.5),
            font=MONO, size=11, color=STONE, align=PP_ALIGN.CENTER)
    textbox(s, project,
            Inches(1), Inches(2.4), Inches(11.33), Inches(1.2),
            font=SERIF, size=48, color=NEAR_BLACK, align=PP_ALIGN.CENTER)
    line(s, Inches(6.17), Inches(3.9), Inches(1), weight=1.5)
    textbox(s, subtitle,
            Inches(1.5), Inches(4.2), Inches(10.33), Inches(0.8),
            font=SANS, size=18, color=OLIVE, align=PP_ALIGN.CENTER)
    textbox(s, f"Registration · {date}",
            Inches(1), Inches(6.5), Inches(11.33), Inches(0.4),
            font=MONO, size=11, color=STONE, align=PP_ALIGN.CENTER)


def content(prs, eyebrow, title, body, page_num=None, items=None):
    s = blank_slide(prs)
    # Eyebrow
    textbox(s, eyebrow.upper(),
            Inches(1.2), Inches(0.5), Inches(10), Inches(0.4),
            font=MONO, size=10, color=STONE)
    # Title
    textbox(s, title,
            Inches(1.2), Inches(1.1), Inches(11), Inches(1.0),
            font=SERIF, size=30, color=NEAR_BLACK)
    line(s, Inches(1.2), Inches(2.2), Inches(11), weight=0.5)
    # Body
    y = Inches(2.5)
    if body:
        textbox(s, body,
                Inches(1.2), y, Inches(11), Inches(1.0),
                font=SANS, size=16, color=OLIVE)
        y += Inches(1.2)
    if items:
        bullet_slide(s, items, Inches(1.2), y, Inches(11), Inches(3.0))
    # Page number
    if page_num is not None:
        textbox(s, f"{page_num:02d}",
                Inches(12), Inches(6.9), Inches(1), Inches(0.3),
                font=MONO, size=10, color=STONE, align=PP_ALIGN.RIGHT)
    return s


def bullet_slide(slide, items, left, top, width, height,
                 font=SANS, size=15, color=DARK_WARM):
    """Add bullet-point text box to an existing slide."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        p.space_before = Pt(0)
        p.level = 0
        run = p.add_run()
        run.text = "• " + item
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tb


def two_col(prs, eyebrow, title,
            left_title, left_items,
            right_title, right_items,
            page_num=None):
    s = blank_slide(prs)
    textbox(s, eyebrow.upper(),
            Inches(1.2), Inches(0.5), Inches(10), Inches(0.4),
            font=MONO, size=10, color=STONE)
    textbox(s, title,
            Inches(1.2), Inches(1.1), Inches(11), Inches(0.8),
            font=SERIF, size=30, color=NEAR_BLACK)
    line(s, Inches(1.2), Inches(2.0), Inches(11), weight=0.5)
    # Left
    textbox(s, left_title,
            Inches(1.2), Inches(2.3), Inches(5), Inches(0.5),
            font=SERIF, size=18, color=OLIVE)
    bullet_slide(s, left_items, Inches(1.2), Inches(2.9), Inches(5), Inches(3.5))
    # Right
    textbox(s, right_title,
            Inches(7.0), Inches(2.3), Inches(5), Inches(0.5),
            font=SERIF, size=18, color=NEAR_BLACK)
    bullet_slide(s, right_items, Inches(7.0), Inches(2.9), Inches(5.2), Inches(3.5))
    if page_num is not None:
        textbox(s, f"{page_num:02d}",
                Inches(12), Inches(6.9), Inches(1), Inches(0.3),
                font=MONO, size=10, color=STONE, align=PP_ALIGN.RIGHT)
    return s


def chapter(prs, number, title):
    s = blank_slide(prs, bg=BRAND)
    textbox(s, f"0{number}",
            Inches(0.8), Inches(0.5), Inches(2), Inches(0.6),
            font=MONO, size=18, color=WHITE)
    textbox(s, title,
            Inches(1), Inches(3.0), Inches(11.33), Inches(1.5),
            font=SERIF, size=48, color=WHITE, align=PP_ALIGN.CENTER)


def metrics(prs, title, metrics_list):
    s = blank_slide(prs)
    textbox(s, title,
            Inches(1.2), Inches(0.8), Inches(11), Inches(0.8),
            font=SERIF, size=28, color=NEAR_BLACK, align=PP_ALIGN.CENTER)
    line(s, Inches(6.17), Inches(1.8), Inches(1))

    n = len(metrics_list)
    cw = Inches(2.6)
    gap = Inches(0.3)
    tw = cw * n + gap * (n - 1)
    start = (SLIDE_W - tw) / 2

    for i, (val, lab) in enumerate(metrics_list):
        x = start + (cw + gap) * i
        textbox(s, val,
                x, Inches(2.5), cw, Inches(1.2),
                font=SERIF, size=48, color=BRAND, align=PP_ALIGN.CENTER)
        textbox(s, lab,
                x, Inches(4.0), cw, Inches(0.6),
                font=SANS, size=12, color=OLIVE, align=PP_ALIGN.CENTER)
    return s


def ending(prs, message, detail, contact):
    s = blank_slide(prs)
    textbox(s, message,
            Inches(1), Inches(2.8), Inches(11.33), Inches(1.2),
            font=SERIF, size=42, color=NEAR_BLACK, align=PP_ALIGN.CENTER)
    textbox(s, detail,
            Inches(1.5), Inches(4.2), Inches(10.33), Inches(0.6),
            font=SANS, size=16, color=OLIVE, align=PP_ALIGN.CENTER)
    line(s, Inches(6.17), Inches(5.0), Inches(1), weight=1.5)
    textbox(s, contact,
            Inches(1), Inches(5.3), Inches(11.33), Inches(0.5),
            font=MONO, size=11, color=STONE, align=PP_ALIGN.CENTER)


def pipeline(prs, eyebrow, title, steps, page_num=None):
    """steps: [(step_title, step_desc), ...]"""
    s = blank_slide(prs)
    textbox(s, eyebrow.upper(),
            Inches(1.2), Inches(0.5), Inches(10), Inches(0.4),
            font=MONO, size=10, color=STONE)
    textbox(s, title,
            Inches(1.2), Inches(1.0), Inches(11), Inches(0.7),
            font=SERIF, size=28, color=NEAR_BLACK)
    line(s, Inches(1.2), Inches(1.85), Inches(11), weight=0.5)

    n = len(steps[:4])
    sw = Inches(11.2 / n)
    for i, (st, sd) in enumerate(steps[:4]):
        x = Inches(1.0) + sw * i
        textbox(s, f"0{i+1}",
                x, Inches(2.2), sw, Inches(0.7),
                font=SERIF, size=36, color=BRAND)
        textbox(s, st,
                x, Inches(3.1), sw - Inches(0.2), Inches(0.6),
                font=SERIF, size=16, color=NEAR_BLACK)
        textbox(s, sd,
                x, Inches(3.8), sw - Inches(0.2), Inches(2.5),
                font=SANS, size=13, color=OLIVE)
    if page_num is not None:
        textbox(s, f"{page_num:02d}",
                Inches(12), Inches(6.9), Inches(1), Inches(0.3),
                font=MONO, size=10, color=STONE, align=PP_ALIGN.RIGHT)
    return s


# ── Main ──────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--out", default="registration-pitch.pptx",
                        help="Output PPTX path")
    args = parser.parse_args()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ── Slide 1: Cover ──
    cover(prs,
          team_name="TEAM NAME TBD",
          project="AI Voice",
          subtitle="On-Device Vietnamese Speech-to-Speech Translation\nfor Industrial Environments",
          date="June 2026")

    # ── Slide 2: Problem ──
    s = blank_slide(prs)
    textbox(s, "PROBLEM",
            Inches(1.2), Inches(0.5), Inches(10), Inches(0.4),
            font=MONO, size=10, color=STONE)
    textbox(s, "Industrial communication fails when workers don't share a spoken language",
            Inches(1.2), Inches(1.0), Inches(11), Inches(1.0),
            font=SERIF, size=28, color=NEAR_BLACK)
    line(s, Inches(1.2), Inches(2.2), Inches(11), weight=0.5)
    bullet_slide(s, [
        "Shift handovers and safety briefings rely on spoken Vietnamese key phrases",
        "Non-Vietnamese-speaking supervisors (technical expats, foreign partners) depend on interpreters",
        "Interpreters introduce latency, errors, and privacy concerns in sensitive settings",
        "Current fallback: paper translations, static phrasebooks, or disconnected cloud apps",
        "$500M+ annual productivity loss in Vietnamese industrial sectors from language barriers",
        "Competing products (Timekettle, Google Translate) require cloud connectivity or specialized HW",
    ], Inches(1.2), Inches(2.5), Inches(11), Inches(3.5), size=15)
    textbox(s, "01",
            Inches(12), Inches(6.9), Inches(1), Inches(0.3),
            font=MONO, size=10, color=STONE, align=PP_ALIGN.RIGHT)

    # ── Slide 3: Solution ──
    s = blank_slide(prs)
    textbox(s, "SOLUTION",
            Inches(1.2), Inches(0.5), Inches(10), Inches(0.4),
            font=MONO, size=10, color=STONE)
    textbox(s, "AI Voice: on-device speech-to-speech translation running entirely on Snapdragon phones",
            Inches(1.2), Inches(1.0), Inches(11), Inches(1.0),
            font=SERIF, size=26, color=NEAR_BLACK)
    line(s, Inches(1.2), Inches(2.2), Inches(11), weight=0.5)
    two_col_blocks = [
        ("Core capability",
         "Worker speaks Vietnamese → phone translates → audio output in English (or vice versa)"),
    ]
    bullet_slide(s, [
        "Worker speaks Vietnamese → phone translates → audio output in target language",
        "No internet required — all inference runs on Qualcomm Snapdragon 8 Gen 2 NPU",
        "Phone-in-pocket deployment via VAD activation + Bluetooth earpiece",
        "Target scenarios: shift handovers, safety briefings, equipment training, incident reporting",
        "Good-enough accuracy prioritized over near-native quality — industrial, not consumer",
    ], Inches(1.2), Inches(2.5), Inches(11), Inches(3.2), size=15)
    textbox(s, "02",
            Inches(12), Inches(6.9), Inches(1), Inches(0.3),
            font=MONO, size=10, color=STONE, align=PP_ALIGN.RIGHT)

    # ── Slide 4: Pipeline ──
    pipeline(prs,
             eyebrow="HOW IT WORKS",
             title="A cascaded pipeline of open-source models optimized for on-device inference",
             steps=[
                 ("Voice Activity\nDetection",
                  "Silero VAD (mit-hr-lab)\n~15ms latency\nFrames audio start/end"),
                 ("Speech\nRecognition",
                  "Whisper Medium (OpenAI)\nvia Qualcomm AI Hub int8\n~400-600ms, 8 languages"),
                 ("Machine\nTranslation",
                  "Quantized 440M-parameter NMT\nCPU-optimized inference\n~TBD ms, VI↔EN/ZH/KO"),
                 ("Text-to-\nSpeech",
                  "MeloTTS-VI custom\nQualcomm Hub extension\n~150-250ms, native VI voice"),
             ],
             page_num=3)

    # ── Slide 5: Why Snapdragon ──
    chapter(prs, "01", "Why Snapdragon 8 Gen 2")

    # ── Slide 6: SD8G2 argument ──
    two_col(prs,
            eyebrow="FORM FACTOR",
            title="The phone workers already own is the best Edge AI device",
            left_title="Wearables (the expensive path)",
            left_items=[
                "Weak battery for sustained inference",
                "No NPU — runs on CPU -> thermal failure",
                "$300-500 additional hardware per worker",
                "Overheats within 5-10 min of continuous load",
                "Unproven durability in industrial settings",
            ],
            right_title="SD8G2 Phone (our path)",
            right_items=[
                "16GB RAM, dedicated Hexagon NPU + Adreno GPU",
                "50M+ devices shipped (Galaxy S23, OnePlus 11, Xiaomi 13)",
                "Thermal sweet spot: TSMC 4nm mature node",
                "Zero additional hardware cost in enterprise deployment",
                "$2-9/device/month via existing MDM (Samsung Knox, VMware)",
            ],
            page_num=5)

    # ── Slide 7: Our edge ──
    s = blank_slide(prs)
    textbox(s, "DIFFERENTIATION",
            Inches(1.2), Inches(0.5), Inches(10), Inches(0.4),
            font=MONO, size=10, color=STONE)
    textbox(s, "Three competitive edges that a tech-burger solution cannot replicate",
            Inches(1.2), Inches(1.0), Inches(11), Inches(1.0),
            font=SERIF, size=26, color=NEAR_BLACK)
    line(s, Inches(1.2), Inches(2.2), Inches(11), weight=0.5)

    edges = [
        ("MeloTTS-VI\nQualcomm Hub Contribution",
         "First Vietnamese TTS in the MeloTTS family on Qualcomm AI Hub. Institutionally aligned with the competition's goal of expanding the Hub catalog. Builds on community checkpoint (nmcuong/MeloTTS-Vietnamese)."),
        ("DeepFilterNet\nTonal Preservation",
         "Leading hypothesis: noise suppression with tonal-language-aware enhancement reduces WER by 15-30% in industrial noise. We plan to validate this empirically — proprietary data no other team will have."),
        ("Thermal Tiers\nGraceful Degradation",
         "Three-tier inference plan (nominal/warm/throttled) profiled on SD8G2. Most teams' demos crash under heat; ours falls back gracefully with staged quality reduction."),
    ]
    y = Inches(2.5)
    for title, desc in edges:
        card(s, Inches(1.2), y, Inches(11), Inches(1.35))
        textbox(s, title,
                Inches(1.5), y + Inches(0.1), Inches(3.5), Inches(1.1),
                font=SERIF, size=15, color=BRAND)
        textbox(s, desc,
                Inches(5.0), y + Inches(0.1), Inches(7.0), Inches(1.1),
                font=SANS, size=12, color=DARK_WARM)
        y += Inches(1.55)
    textbox(s, "06",
            Inches(12), Inches(6.9), Inches(1), Inches(0.3),
            font=MONO, size=10, color=STONE, align=PP_ALIGN.RIGHT)

    # ── Slide 8: MeloTTS-VI detail ──
    two_col(prs,
            eyebrow="INNOVATION",
            title="MeloTTS-VI: the first Vietnamese TTS on Qualcomm AI Hub",
            left_title="Community checkpoint exists",
            left_items=[
                "nmcuong/MeloTTS-Vietnamese on HuggingFace",
                "MIT licensed — no IP conflict",
                "25-hour Infore dataset (acknowledged quality limits)",
                "PhoBERT-based encoder (architectural divergence)",
                "Saves months vs. training from scratch",
            ],
            right_title="Our plan to productionize",
            right_items=[
                "Fine-tune with industrial-noise-domain data",
                "Adapt for Qualcomm AI Hub Workbench export",
                "ONNX quantize with custom calibration",
                "Validate MOS with native speaker panel",
                "Ship back to Hub as first-party Vietnamese TTS",
            ],
            page_num=7)

    # ── Slide 9: Market gap ──
    two_col(prs,
            eyebrow="MARKET",
            title="No enterprise product solves Vietnamese S2ST offline in noise",
            left_title="What exists today",
            left_items=[
                "Google Translate: cloud-dependent, no S2ST pipeline",
                "Timekettle WT2 (hardware): $299, cloud hybrid",
                "RTranslator (open source): Whisper+NLLB, unoptimized",
                "Quicksilver, iFLYTEK: Chinese-focused, cloud-heavy",
            ],
            right_title="What's missing (our slot)",
            right_items=[
                "Complete offline S2ST on common phone hardware",
                "Vietnamese as primary language (not EN-first port)",
                "Industrial noise robustness with empirical validation",
                "Qualcomm AI Hub model contribution (ecosystem play)",
                "Thermal-aware deployment for sustained industrial use",
            ],
            page_num=8)

    # ── Slide 10: Business viability ──
    s = blank_slide(prs)
    textbox(s, "BUSINESS VIABILITY",
            Inches(1.2), Inches(0.5), Inches(10), Inches(0.4),
            font=MONO, size=10, color=STONE)
    textbox(s, "Enterprise deployability without specialized hardware investment",
            Inches(1.2), Inches(1.0), Inches(11), Inches(0.8),
            font=SERIF, size=28, color=NEAR_BLACK)
    line(s, Inches(1.2), Inches(2.0), Inches(11), weight=0.5)

    metrics(prs, "", [
        ("50M+", "SD8G2 devices\nshipped globally"),
        ("$2-9", "MDM cost\n/device/month"),
        ("35-50%", "Refurbished\nprice discount"),
        ("3+ yrs", "Enterprise\nsupport window"),
    ])
    bullet_slide(s, [
        "Zero incremental hardware — deployable via existing Samsung Knox / VMware MDM infrastructure",
        "Target: footwear, garment, electronics factories with Vietnamese workforce + foreign management",
        "Growing refurbished SD8G2 market reduces procurement cost for enterprise pilot programs",
        "Competition prize includes partnership for product development — de-risks go-to-market",
    ], Inches(1.2), Inches(4.5), Inches(11), Inches(2.5), size=14)
    textbox(s, "09",
            Inches(12), Inches(6.9), Inches(1), Inches(0.3),
            font=MONO, size=10, color=STONE, align=PP_ALIGN.RIGHT)

    # ── Slide 11: Roadmap ──
    pipeline(prs,
             eyebrow="ROADMAP",
             title="Four phases from prototype to field-ready product",
             steps=[
                 ("Phase 1\nJun-Jul 2026",
                  "Prototype pipeline on SD8G2\nFocused on VI-EN pair\nLatency budget validation"),
                 ("Phase 2\nJul-Aug 2026",
                  "Noise robustness tuning\nThermal profiling\nMeloTTS fine-tuning"),
                 ("Phase 3\nAug-Sep 2026",
                  "Field testing at partner site\nMulti-language expansion\nMOS panel validation"),
                 ("Phase 4\nSep-Nov 2026",
                  "Grand Finale presentation\nBusiness case finalization\nHub distribution launch"),
             ],
             page_num=10)

    # ── Slide 12: Closing ──
    ending(prs,
           message="From prototype to product.",
           detail="On-device Vietnamese speech translation for the industrial workforce.",
           contact="TEAM NAME · TEAM EMAIL · GITHUB REPO")

    # Save
    prs.save(args.out)
    print(f"OK: Saved {args.out}")


if __name__ == "__main__":
    main()
